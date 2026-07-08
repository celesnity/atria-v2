"""Multi-format text extraction (PDF / DOCX / PPTX / text / images)."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

_SCRIPTS = Path(__file__).resolve().parent.parent / "modules" / "ai_workspace" / "scripts"
if str(_SCRIPTS) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS))

import extract  # noqa: E402


def test_txt_and_md(tmp_path):
    f = tmp_path / "note.md"
    f.write_text("# Tiêu đề\nNội dung nội bộ.", encoding="utf-8")
    assert "Nội dung nội bộ" in extract.extract_text(str(f))


def test_docx(tmp_path):
    docx = pytest.importorskip("docx")
    path = tmp_path / "policy.docx"
    doc = docx.Document()
    doc.add_paragraph("Chính sách nghỉ phép công ty")
    doc.add_paragraph("Nhân viên có 15 ngày phép.")
    doc.save(str(path))
    text = extract.extract_text(str(path))
    assert "Chính sách nghỉ phép" in text and "15 ngày phép" in text


def test_pptx(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    path = tmp_path / "deck.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Kế hoạch quý 4"
    prs.save(str(path))
    text = extract.extract_text(str(path))
    assert "Kế hoạch quý 4" in text


def test_pdf(tmp_path):
    pypdf = pytest.importorskip("pypdf")
    path = tmp_path / "blank.pdf"
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with open(path, "wb") as fh:
        writer.write(fh)
    # A blank PDF has no extractable text — must not crash, returns a string.
    assert isinstance(extract.extract_text(str(path)), str)


def test_unsupported_format_raises(tmp_path):
    f = tmp_path / "archive.zip"
    f.write_bytes(b"PK\x03\x04")
    with pytest.raises(extract.ExtractionError):
        extract.extract_text(str(f))


def test_is_supported():
    assert extract.is_supported("a.pdf") and extract.is_supported("b.DOCX")
    assert extract.is_supported("c.pptx") and extract.is_supported("d.png")
    assert not extract.is_supported("e.zip")


def test_image_ocr_is_graceful(tmp_path):
    # Without a tesseract backend this returns "" rather than raising.
    Image = pytest.importorskip("PIL.Image")
    p = tmp_path / "pic.png"
    Image.new("RGB", (20, 20), "white").save(str(p))
    assert isinstance(extract.extract_text(str(p)), str)
