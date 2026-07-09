"""Extract plain text from uploaded documents.

Pure-Python extractors for PDF / DOCX / PPTX / plain text. Image OCR is
best-effort: it needs an OCR backend (``pytesseract`` + the tesseract binary);
if that is unavailable the image is still accepted but yields no text.

The extracted text feeds two things: viewing an uploaded document in the
workspace, and (later) embedding it for AI search.
"""
from __future__ import annotations

from pathlib import Path

TEXT_EXT = {".md", ".txt", ".csv", ".json", ".log"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}
SUPPORTED_EXT = TEXT_EXT | PDF_EXT | DOCX_EXT | PPTX_EXT | IMAGE_EXT


class ExtractionError(Exception):
    """Raised when a file cannot be parsed into text."""


def _ext(name: str) -> str:
    return Path(name).suffix.lower()


def is_supported(name: str) -> bool:
    """Whether the filename's extension is a format we can ingest."""
    return _ext(name) in SUPPORTED_EXT


def extract_text(path: str, filename: str | None = None) -> str:
    """Extract text from ``path``. ``filename`` drives format detection.

    Returns the extracted text (possibly empty, e.g. a scanned image with no OCR
    backend). Raises :class:`ExtractionError` for unsupported formats.
    """
    name = filename or Path(path).name
    ext = _ext(name)
    if ext in TEXT_EXT:
        return Path(path).read_text(encoding="utf-8", errors="replace").strip()
    if ext in PDF_EXT:
        return _extract_pdf(path)
    if ext in DOCX_EXT:
        return _extract_docx(path)
    if ext in PPTX_EXT:
        return _extract_pptx(path)
    if ext in IMAGE_EXT:
        return _extract_image(path)
    raise ExtractionError(f"unsupported format: {ext or name!r}")


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = [(page.extract_text() or "") for page in reader.pages]
    return "\n\n".join(p for p in pages if p).strip()


def _extract_docx(path: str) -> str:
    import docx

    document = docx.Document(path)
    parts: list[str] = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(p for p in parts if p and p.strip()).strip()


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
    return "\n".join(parts).strip()


def _extract_image(path: str) -> str:
    """OCR an image, or return "" if no OCR backend is installed."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""
    img = Image.open(path)
    # Prefer Vietnamese+English; fall back to the default language pack, then to
    # "" if the tesseract binary is missing entirely.
    for lang in ("vie+eng", None):
        try:
            return pytesseract.image_to_string(img, lang=lang).strip()
        except Exception:  # noqa: BLE001 - lang pack or binary missing
            continue
    return ""
